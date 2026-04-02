from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────
DARK_BG    = RGBColor(0x0F, 0x0F, 0x1A)
MID_BG     = RGBColor(0x16, 0x21, 0x3E)
BLUE_DARK  = RGBColor(0x0F, 0x34, 0x60)
MS_BLUE    = RGBColor(0x00, 0x78, 0xD4)
G_BLUE     = RGBColor(0x1A, 0x73, 0xE8)
ACCENT     = RGBColor(0x60, 0xA5, 0xFA)
ACCENT2    = RGBColor(0xA7, 0x8B, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM  = RGBColor(0xCC, 0xCC, 0xDD)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
RED_SOFT   = RGBColor(0xE5, 0x73, 0x73)

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ─────────────────────────────────────────────────

def add_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  font_size=16, color=WHITE_DIM, bold_first=False, spacing=1.15):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
        run.font.name = "Calibri"

def divider(slide, l, t, w=0.5, color=ACCENT):
    add_rect(slide, l, t, w, 0.04, color)

# ════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1, DARK_BG)
add_rect(s1, 0, 0, 13.33, 7.5, MID_BG)          # subtle overlay
add_rect(s1, 0, 0, 5.5, 7.5, BLUE_DARK)          # left accent panel
add_rect(s1, 0, 0, 0.12, 7.5, ACCENT)            # left edge stripe

add_textbox(s1, "GROUP 9 PRESENTATION", 0.5, 0.5, 5, 0.4,
            font_size=9, color=ACCENT, bold=True)
add_textbox(s1, "Microsoft vs\nGoogle Tools\nfor HR Optimisation",
            0.5, 1.1, 5.2, 3.5, font_size=36, bold=True, color=WHITE)
add_textbox(s1, "A comparative analysis of collaborative productivity\nplatforms and their impact on modern HR operations.",
            0.5, 4.5, 5.2, 1.5, font_size=13, color=WHITE_DIM)

# right side decorative circles (simulated with rounded rects)
add_rect(s1, 7.5, -1.5, 5, 5, RGBColor(0x1E, 0x2A, 0x4A))
add_rect(s1, 9,    4,   4, 4, RGBColor(0x12, 0x1E, 0x35))
add_textbox(s1, "HR Technology  ·  Group 9", 6.5, 6.8, 6, 0.4,
            font_size=10, color=RGBColor(0x55,0x66,0x88), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2, DARK_BG)
add_rect(s2, 0, 0, 13.33, 1.4, MID_BG)
add_rect(s2, 0, 0, 13.33, 0.06, ACCENT)

add_textbox(s2, "OVERVIEW", 0.6, 0.18, 4, 0.3, font_size=9, color=ACCENT, bold=True)
add_textbox(s2, "Agenda", 0.6, 0.45, 6, 0.7, font_size=28, bold=True, color=WHITE)

items = [
    ("01", "Definition",        "Collaborative tools & their role in HR"),
    ("02", "Microsoft 365",     "Key tools and HR applications"),
    ("03", "Google Workspace",  "Key tools and HR applications"),
    ("04", "Comparison",        "Integration, Security, Data Management"),
    ("05", "Conclusion",        "Which platform suits HR best?"),
]

for i, (num, title, desc) in enumerate(items):
    row_t = 1.6 + i * 1.0
    add_rect(s2, 0.6, row_t, 0.55, 0.55, ACCENT)
    add_textbox(s2, num, 0.6, row_t, 0.55, 0.55,
                font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_textbox(s2, title, 1.3, row_t - 0.02, 3, 0.35, font_size=13, bold=True, color=WHITE)
    add_textbox(s2, desc,  1.3, row_t + 0.3,  8, 0.3,  font_size=10, color=WHITE_DIM)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Definition
# ════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3, RGBColor(0x1A, 0x1A, 0x2E))
add_rect(s3, 0, 0, 13.33, 1.4, RGBColor(0x2D, 0x1B, 0x69))
add_rect(s3, 0, 0, 13.33, 0.06, ACCENT)

add_textbox(s3, "SECTION 01", 0.6, 0.18, 4, 0.3, font_size=9, color=ACCENT, bold=True)
add_textbox(s3, "Collaborative Productivity Tools", 0.6, 0.45, 10, 0.7,
            font_size=26, bold=True, color=WHITE)

# Definition box
add_rect(s3, 0.6, 1.6, 5.7, 2.1, RGBColor(0x25, 0x25, 0x45))
add_rect(s3, 0.6, 1.6, 0.06, 2.1, ACCENT)
add_textbox(s3, "DEFINITION", 0.8, 1.68, 5, 0.3, font_size=8, bold=True, color=ACCENT)
add_textbox(s3,
    "Digital platforms that enable teams to communicate, share information, "
    "manage tasks, and work together in real time — regardless of physical location.",
    0.8, 2.0, 5.3, 1.5, font_size=12, color=WHITE_DIM)

# HR Role box
add_rect(s3, 0.6, 3.85, 5.7, 2.1, RGBColor(0x25, 0x25, 0x45))
add_rect(s3, 0.6, 3.85, 0.06, 2.1, ACCENT2)
add_textbox(s3, "ROLE IN HR OPERATIONS", 0.8, 3.93, 5, 0.3, font_size=8, bold=True, color=ACCENT2)
add_textbox(s3,
    "Streamline recruitment, onboarding, performance management, payroll "
    "coordination, and employee communication — reducing manual effort and "
    "improving organisational efficiency.",
    0.8, 4.25, 5.3, 1.5, font_size=12, color=WHITE_DIM)

# Chips on right
chips = ["Recruitment", "Onboarding", "Performance Reviews", "Payroll", "Communication", "Training"]
for i, chip in enumerate(chips):
    cx = 7.0 + (i % 2) * 2.9
    cy = 2.0 + (i // 2) * 0.75
    add_rect(s3, cx, cy, 2.6, 0.45, RGBColor(0x1E, 0x2A, 0x50))
    add_textbox(s3, chip, cx + 0.1, cy + 0.04, 2.4, 0.38,
                font_size=11, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Microsoft 365
# ════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4, RGBColor(0x00, 0x50, 0x9A))
add_rect(s4, 0, 0, 13.33, 1.4, RGBColor(0x00, 0x3F, 0x7A))
add_rect(s4, 0, 0, 13.33, 0.06, WHITE)

add_textbox(s4, "SECTION 02", 0.6, 0.18, 4, 0.3, font_size=9,
            color=RGBColor(0xBB, 0xDD, 0xFF), bold=True)
add_textbox(s4, "Microsoft 365 for HR", 0.6, 0.45, 8, 0.7,
            font_size=26, bold=True, color=WHITE)

# Left card
add_rect(s4, 0.6, 1.6, 5.7, 5.3, RGBColor(0x00, 0x3F, 0x7A))
add_textbox(s4, "💼  Core HR Tools", 0.8, 1.75, 5, 0.45, font_size=13, bold=True, color=WHITE)
tools = [
    "Teams — video calls, chat & HR meetings",
    "SharePoint — document & policy management",
    "Outlook — professional email & scheduling",
    "Excel — payroll, analytics & reporting",
    "Forms — employee surveys & feedback",
    "Power Automate — automated HR workflows",
]
for i, t in enumerate(tools):
    add_rect(s4, 0.8, 2.3 + i*0.72, 0.06, 0.45, ACCENT)
    add_textbox(s4, t, 1.0, 2.3 + i*0.72, 5.1, 0.45, font_size=11, color=WHITE_DIM)

# Right card
add_rect(s4, 6.9, 1.6, 5.8, 5.3, RGBColor(0x00, 0x3F, 0x7A))
add_textbox(s4, "🚀  HR Benefits", 7.1, 1.75, 5, 0.45, font_size=13, bold=True, color=WHITE)
benefits = [
    "Deep integration with Azure AD & Windows",
    "Viva Insights for employee wellbeing",
    "Compliance Manager for regulated industries",
    "Enterprise-grade security & audit logs",
    "Power BI dashboards for HR analytics",
    "Granular data loss prevention (DLP) policies",
]
for i, b in enumerate(benefits):
    add_rect(s4, 7.1, 2.3 + i*0.72, 0.06, 0.45, GOLD)
    add_textbox(s4, b, 7.3, 2.3 + i*0.72, 5.2, 0.45, font_size=11, color=WHITE_DIM)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Google Workspace
# ════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_bg(s5, RGBColor(0x0D, 0x47, 0xA1))
add_rect(s5, 0, 0, 13.33, 1.4, RGBColor(0x0A, 0x35, 0x7A))
add_rect(s5, 0, 0, 13.33, 0.06, WHITE)

add_textbox(s5, "SECTION 03", 0.6, 0.18, 4, 0.3, font_size=9,
            color=RGBColor(0xBB, 0xDD, 0xFF), bold=True)
add_textbox(s5, "Google Workspace for HR", 0.6, 0.45, 8, 0.7,
            font_size=26, bold=True, color=WHITE)

add_rect(s5, 0.6, 1.6, 5.7, 5.3, RGBColor(0x0A, 0x35, 0x7A))
add_textbox(s5, "🌐  Core HR Tools", 0.8, 1.75, 5, 0.45, font_size=13, bold=True, color=WHITE)
gtools = [
    "Meet — video interviews & team calls",
    "Drive — cloud document storage & sharing",
    "Gmail — communication & scheduling",
    "Sheets — data tracking & HR analytics",
    "Forms — job applications & surveys",
    "AppSheet — no-code HR app builder",
]
for i, t in enumerate(gtools):
    add_rect(s5, 0.8, 2.3 + i*0.72, 0.06, 0.45, ACCENT)
    add_textbox(s5, t, 1.0, 2.3 + i*0.72, 5.1, 0.45, font_size=11, color=WHITE_DIM)

add_rect(s5, 6.9, 1.6, 5.8, 5.3, RGBColor(0x0A, 0x35, 0x7A))
add_textbox(s5, "✨  HR Benefits", 7.1, 1.75, 5, 0.45, font_size=13, bold=True, color=WHITE)
gbenefits = [
    "Real-time collaboration on all documents",
    "Simple, intuitive UI — low training cost",
    "Strong mobile experience for remote HR",
    "Google Workspace Marketplace integrations",
    "AI-powered search across all HR files",
    "Zero-trust BeyondCorp security model",
]
for i, b in enumerate(gbenefits):
    add_rect(s5, 7.1, 2.3 + i*0.72, 0.06, 0.45, GOLD)
    add_textbox(s5, b, 7.3, 2.3 + i*0.72, 5.2, 0.45, font_size=11, color=WHITE_DIM)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — Comparison Table
# ════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_bg(s6, RGBColor(0x0D, 0x0D, 0x1A))
add_rect(s6, 0, 0, 13.33, 1.4, RGBColor(0x15, 0x05, 0x33))
add_rect(s6, 0, 0, 13.33, 0.06, ACCENT)

add_textbox(s6, "SECTION 04", 0.6, 0.18, 4, 0.3, font_size=9, color=ACCENT, bold=True)
add_textbox(s6, "Head-to-Head Comparison", 0.6, 0.45, 10, 0.7,
            font_size=26, bold=True, color=WHITE)

# Table header
add_rect(s6, 0.5, 1.55, 2.2, 0.5, RGBColor(0x22, 0x22, 0x44))
add_rect(s6, 2.75, 1.55, 4.7, 0.5, MS_BLUE)
add_rect(s6, 7.5, 1.55, 5.2, 0.5, G_BLUE)
add_textbox(s6, "Criteria",        0.55, 1.6, 2.1, 0.4, font_size=11, bold=True, color=ACCENT)
add_textbox(s6, "🏢  Microsoft 365", 2.8, 1.6, 4.5, 0.4, font_size=11, bold=True, color=WHITE)
add_textbox(s6, "🌐  Google Workspace", 7.55, 1.6, 5.0, 0.4, font_size=11, bold=True, color=WHITE)

rows = [
    ("Integration",
     "Deep integration with Azure AD, SAP, Dynamics 365 — ideal for large enterprises.",
     "Strong with Google Cloud & third-party apps via Marketplace — flexible for modern stacks."),
    ("Security",
     "Advanced Threat Protection, Compliance Manager, GDPR tools, enterprise audit logs.",
     "Zero-trust BeyondCorp model, 2FA, Vault for eDiscovery — strong but simpler controls."),
    ("Data\nManagement",
     "SharePoint + OneDrive, granular permissions, robust data loss prevention (DLP) policies.",
     "Google Drive with shared drives, AI-powered search, simpler permission model."),
]

row_colors = [RGBColor(0x18,0x18,0x30), RGBColor(0x1E,0x1E,0x38), RGBColor(0x18,0x18,0x30)]

for i, (crit, ms, goog) in enumerate(rows):
    ry = 2.15 + i * 1.55
    add_rect(s6, 0.5,  ry, 2.2,  1.45, row_colors[i])
    add_rect(s6, 2.75, ry, 4.7,  1.45, row_colors[i])
    add_rect(s6, 7.5,  ry, 5.2,  1.45, row_colors[i])
    add_rect(s6, 0.5,  ry, 0.06, 1.45, ACCENT)
    add_textbox(s6, crit, 0.65, ry+0.1, 2.0, 1.2, font_size=12, bold=True, color=ACCENT)
    add_textbox(s6, ms,   2.85, ry+0.08, 4.4, 1.3, font_size=10, color=WHITE_DIM)
    add_textbox(s6, goog, 7.6,  ry+0.08, 4.9, 1.3, font_size=10, color=WHITE_DIM)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Verdict
# ════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_bg(s7, RGBColor(0x1A, 0x1A, 0x2E))
add_rect(s7, 0, 0, 13.33, 1.4, RGBColor(0x2D, 0x1B, 0x69))
add_rect(s7, 0, 0, 13.33, 0.06, ACCENT)

add_textbox(s7, "SECTION 04 — DEEP DIVE", 0.6, 0.18, 6, 0.3, font_size=9, color=ACCENT, bold=True)
add_textbox(s7, "Which Wins in Each Area?", 0.6, 0.45, 10, 0.7,
            font_size=26, bold=True, color=WHITE)

verdicts = [
    ("🔗", "Integration",
     "Microsoft 365 wins for large enterprises with existing Microsoft ecosystems.\n"
     "Google wins for cloud-native, agile organisations."),
    ("🔒", "Security",
     "Microsoft offers more granular compliance controls for regulated industries.\n"
     "Google provides a simpler, modern zero-trust approach."),
    ("🗄️", "Data Management",
     "Microsoft's DLP and SharePoint give HR teams more control and governance.\n"
     "Google Drive is easier to use but less configurable at scale."),
]

for i, (icon, title, desc) in enumerate(verdicts):
    cx = 0.6 + i * 4.2
    add_rect(s7, cx, 1.7, 3.9, 4.8, RGBColor(0x25, 0x20, 0x50))
    add_rect(s7, cx, 1.7, 3.9, 0.06, ACCENT)
    add_textbox(s7, icon,  cx+0.15, 1.85, 3.6, 0.7, font_size=28, align=PP_ALIGN.CENTER)
    add_textbox(s7, title, cx+0.15, 2.65, 3.6, 0.5, font_size=14, bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s7, desc,  cx+0.2,  3.25, 3.5, 2.8, font_size=11, color=WHITE_DIM)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — Conclusion
# ════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
add_bg(s8, RGBColor(0x0F, 0x20, 0x27))
add_rect(s8, 0, 0, 13.33, 0.06, ACCENT)
add_rect(s8, 0, 7.44, 13.33, 0.06, ACCENT)

add_textbox(s8, "CONCLUSION", 0.6, 0.25, 6, 0.35, font_size=9, color=ACCENT, bold=True)
add_textbox(s8, "🎯", 5.9, 0.6, 1.5, 1.0, font_size=40, align=PP_ALIGN.CENTER)
add_textbox(s8,
    "Both platforms optimise HR —\nthe right choice depends on context.",
    1.5, 1.7, 10.3, 1.4, font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider(s8, 6.1, 3.25, 1.1)

add_textbox(s8,
    "Large enterprises with complex compliance needs lean toward Microsoft 365.\n"
    "Agile, remote-first organisations benefit more from Google Workspace.\n"
    "HR teams should evaluate based on existing infrastructure, team size, and security requirements.",
    1.2, 3.5, 10.9, 2.0, font_size=13, color=WHITE_DIM, align=PP_ALIGN.CENTER)

# chips row
chip_labels = ["Thank You", "Group 9", "Questions Welcome"]
for i, label in enumerate(chip_labels):
    cx = 3.8 + i * 2.0
    add_rect(s8, cx, 5.8, 1.7, 0.45, RGBColor(0x1E, 0x3A, 0x50))
    add_textbox(s8, label, cx+0.05, 5.83, 1.6, 0.38,
                font_size=10, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)

# ── Save ────────────────────────────────────────────────────
prs.save("HR_Optimisation_Group9.pptx")
print("Saved: HR_Optimisation_Group9.pptx")
